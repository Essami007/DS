#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_encg_presentation():
    """Create a comprehensive presentation about ENCG Settat"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "ENCG SETTAT"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "École Nationale de Commerce et de Gestion"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Présentation Générale
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Présentation Générale"
    title2.text_frame.paragraphs[0].font.size = Pt(44)
    title2.text_frame.paragraphs[0].font.bold = True
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "L'ENCG Settat"
    
    p1 = tf2.add_paragraph()
    p1.text = "Établissement public d'enseignement supérieur"
    p1.level = 1
    
    p2 = tf2.add_paragraph()
    p2.text = "Créée en 2004"
    p2.level = 1
    
    p3 = tf2.add_paragraph()
    p3.text = "Fait partie du réseau des ENCG du Maroc"
    p3.level = 1
    
    p4 = tf2.add_paragraph()
    p4.text = "Formation en management et gestion"
    p4.level = 1
    
    p5 = tf2.add_paragraph()
    p5.text = "Située à Settat, région de Casablanca-Settat"
    p5.level = 1
    
    for paragraph in tf2.paragraphs:
        paragraph.font.size = Pt(24)
    
    # Slide 3: Mission et Vision
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Mission et Vision"
    title3.text_frame.paragraphs[0].font.size = Pt(44)
    title3.text_frame.paragraphs[0].font.bold = True
    title3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Mission"
    
    p1 = tf3.add_paragraph()
    p1.text = "Former des cadres compétents en management"
    p1.level = 1
    
    p2 = tf3.add_paragraph()
    p2.text = "Développer les compétences entrepreneuriales"
    p2.level = 1
    
    p3 = tf3.add_paragraph()
    p3.text = "Promouvoir la recherche scientifique"
    p3.level = 1
    
    p4 = tf3.add_paragraph()
    p4.text = ""
    p4.level = 0
    
    p5 = tf3.add_paragraph()
    p5.text = "Vision"
    p5.level = 0
    
    p6 = tf3.add_paragraph()
    p6.text = "Être une école de référence en Afrique"
    p6.level = 1
    
    p7 = tf3.add_paragraph()
    p7.text = "Excellence académique et professionnelle"
    p7.level = 1
    
    for paragraph in tf3.paragraphs:
        paragraph.font.size = Pt(22)
    
    # Slide 4: Formations Proposées
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Formations Proposées"
    title4.text_frame.paragraphs[0].font.size = Pt(44)
    title4.text_frame.paragraphs[0].font.bold = True
    title4.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "Cycle Fondamental (3 ans)"
    
    p1 = tf4.add_paragraph()
    p1.text = "Formation générale en management"
    p1.level = 1
    
    p2 = tf4.add_paragraph()
    p2.text = "Diplôme de Licence"
    p2.level = 1
    
    p3 = tf4.add_paragraph()
    p3.text = ""
    p3.level = 0
    
    p4 = tf4.add_paragraph()
    p4.text = "Cycle de Spécialisation (2 ans)"
    p4.level = 0
    
    p5 = tf4.add_paragraph()
    p5.text = "Finance et Comptabilité"
    p5.level = 1
    
    p6 = tf4.add_paragraph()
    p6.text = "Marketing et Action Commerciale"
    p6.level = 1
    
    p7 = tf4.add_paragraph()
    p7.text = "Management des Ressources Humaines"
    p7.level = 1
    
    p8 = tf4.add_paragraph()
    p8.text = "Gestion des Systèmes d'Information"
    p8.level = 1
    
    for paragraph in tf4.paragraphs:
        paragraph.font.size = Pt(22)
    
    # Slide 5: Conditions d'Admission
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Conditions d'Admission"
    title5.text_frame.paragraphs[0].font.size = Pt(44)
    title5.text_frame.paragraphs[0].font.bold = True
    title5.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "Cycle Fondamental"
    
    p1 = tf5.add_paragraph()
    p1.text = "Baccalauréat (toutes filières)"
    p1.level = 1
    
    p2 = tf5.add_paragraph()
    p2.text = "Concours national d'accès"
    p2.level = 1
    
    p3 = tf5.add_paragraph()
    p3.text = "Épreuves écrites et orales"
    p3.level = 1
    
    p4 = tf5.add_paragraph()
    p4.text = ""
    p4.level = 0
    
    p5 = tf5.add_paragraph()
    p5.text = "Cycle de Spécialisation"
    p5.level = 0
    
    p6 = tf5.add_paragraph()
    p6.text = "Validation du cycle fondamental"
    p6.level = 1
    
    p7 = tf5.add_paragraph()
    p7.text = "Choix de spécialisation selon mérite"
    p7.level = 1
    
    for paragraph in tf5.paragraphs:
        paragraph.font.size = Pt(24)
    
    # Slide 6: Vie Étudiante
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "Vie Étudiante et Activités"
    title6.text_frame.paragraphs[0].font.size = Pt(44)
    title6.text_frame.paragraphs[0].font.bold = True
    title6.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "Clubs et Associations"
    
    p1 = tf6.add_paragraph()
    p1.text = "Club entrepreneuriat"
    p1.level = 1
    
    p2 = tf6.add_paragraph()
    p2.text = "Club culturel et artistique"
    p2.level = 1
    
    p3 = tf6.add_paragraph()
    p3.text = "Club sportif"
    p3.level = 1
    
    p4 = tf6.add_paragraph()
    p4.text = "Événements et séminaires"
    p4.level = 1
    
    p5 = tf6.add_paragraph()
    p5.text = "Stages en entreprise"
    p5.level = 1
    
    p6 = tf6.add_paragraph()
    p6.text = "Projets de fin d'études"
    p6.level = 1
    
    for paragraph in tf6.paragraphs:
        paragraph.font.size = Pt(24)
    
    # Slide 7: Partenariats et Débouchés
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "Partenariats et Débouchés"
    title7.text_frame.paragraphs[0].font.size = Pt(44)
    title7.text_frame.paragraphs[0].font.bold = True
    title7.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "Partenariats"
    
    p1 = tf7.add_paragraph()
    p1.text = "Entreprises nationales et internationales"
    p1.level = 1
    
    p2 = tf7.add_paragraph()
    p2.text = "Universités étrangères"
    p2.level = 1
    
    p3 = tf7.add_paragraph()
    p3.text = "Programmes d'échange"
    p3.level = 1
    
    p4 = tf7.add_paragraph()
    p4.text = ""
    p4.level = 0
    
    p5 = tf7.add_paragraph()
    p5.text = "Débouchés Professionnels"
    p5.level = 0
    
    p6 = tf7.add_paragraph()
    p6.text = "Cadres en entreprises"
    p6.level = 1
    
    p7 = tf7.add_paragraph()
    p7.text = "Consultants"
    p7.level = 1
    
    p8 = tf7.add_paragraph()
    p8.text = "Entrepreneurs"
    p8.level = 1
    
    p9 = tf7.add_paragraph()
    p9.text = "Poursuite d'études (Doctorat, MBA)"
    p9.level = 1
    
    for paragraph in tf7.paragraphs:
        paragraph.font.size = Pt(22)
    
    # Slide 8: Infrastructures
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    title8.text = "Infrastructures et Équipements"
    title8.text_frame.paragraphs[0].font.size = Pt(44)
    title8.text_frame.paragraphs[0].font.bold = True
    title8.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.text = "Campus moderne et équipé"
    
    p1 = tf8.add_paragraph()
    p1.text = "Salles de cours climatisées"
    p1.level = 1
    
    p2 = tf8.add_paragraph()
    p2.text = "Laboratoires informatiques"
    p2.level = 1
    
    p3 = tf8.add_paragraph()
    p3.text = "Bibliothèque riche et moderne"
    p3.level = 1
    
    p4 = tf8.add_paragraph()
    p4.text = "Espaces de coworking"
    p4.level = 1
    
    p5 = tf8.add_paragraph()
    p5.text = "Connexion WiFi haut débit"
    p5.level = 1
    
    p6 = tf8.add_paragraph()
    p6.text = "Installations sportives"
    p6.level = 1
    
    p7 = tf8.add_paragraph()
    p7.text = "Cafétéria et espaces de détente"
    p7.level = 1
    
    for paragraph in tf8.paragraphs:
        paragraph.font.size = Pt(24)
    
    # Slide 9: Chiffres Clés
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    title9.text = "Chiffres Clés"
    title9.text_frame.paragraphs[0].font.size = Pt(44)
    title9.text_frame.paragraphs[0].font.bold = True
    title9.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.text = "Plus de 1000 étudiants"
    
    p1 = tf9.add_paragraph()
    p1.text = ""
    p1.level = 0
    
    p2 = tf9.add_paragraph()
    p2.text = "Corps professoral qualifié (50+ enseignants)"
    p2.level = 0
    
    p3 = tf9.add_paragraph()
    p3.text = ""
    p3.level = 0
    
    p4 = tf9.add_paragraph()
    p4.text = "Taux d'insertion professionnelle élevé"
    p4.level = 0
    
    p5 = tf9.add_paragraph()
    p5.text = ""
    p5.level = 0
    
    p6 = tf9.add_paragraph()
    p6.text = "Réseau de plus de 3000 lauréats"
    p6.level = 0
    
    p7 = tf9.add_paragraph()
    p7.text = ""
    p7.level = 0
    
    p8 = tf9.add_paragraph()
    p8.text = "Nombreux partenariats nationaux et internationaux"
    p8.level = 0
    
    for paragraph in tf9.paragraphs:
        paragraph.font.size = Pt(26)
        paragraph.font.bold = True
    
    # Slide 10: Contact et Conclusion
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background10 = slide10.background
    fill10 = background10.fill
    fill10.solid()
    fill10.fore_color.rgb = RGBColor(0, 51, 102)
    
    # Title
    title_box10 = slide10.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame10 = title_box10.text_frame
    title_frame10.text = "Merci de votre attention !"
    title_para10 = title_frame10.paragraphs[0]
    title_para10.font.size = Pt(48)
    title_para10.font.bold = True
    title_para10.font.color.rgb = RGBColor(255, 255, 255)
    title_para10.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide10.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(3))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    
    contact_frame.text = "Contact"
    p1 = contact_frame.add_paragraph()
    p1.text = ""
    p2 = contact_frame.add_paragraph()
    p2.text = "📍 ENCG Settat"
    p3 = contact_frame.add_paragraph()
    p3.text = "Km 3, Route de Casablanca, Settat"
    p4 = contact_frame.add_paragraph()
    p4.text = ""
    p5 = contact_frame.add_paragraph()
    p5.text = "🌐 www.encg-settat.ma"
    p6 = contact_frame.add_paragraph()
    p6.text = ""
    p7 = contact_frame.add_paragraph()
    p7.text = "📧 contact@encg-settat.ma"
    
    for paragraph in contact_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('/vercel/sandbox/ENCG_SETTAT_Presentation.pptx')
    print("✅ Présentation créée avec succès : ENCG_SETTAT_Presentation.pptx")

if __name__ == "__main__":
    create_encg_presentation()
